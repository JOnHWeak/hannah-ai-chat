import glob
import os
from typing import Optional, Tuple

from database import SessionLocal
from models import KnowledgeBase


def extract_text_from_ppt(path: str) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise RuntimeError(
            "Missing dependency python-pptx. Install with: pip install python-pptx"
        ) from exc

    presentation = Presentation(path)
    texts = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs).strip()
                    if line:
                        texts.append(line)
    return "\n".join(texts)


def import_ppt_dir(
    directory: str,
    category: str,
    created_by: str = "ppt_import",
) -> int:
    session = SessionLocal()
    count = 0
    try:
        patterns = ["**/*.ppt", "**/*.pptx"]
        for pattern in patterns:
            for path in glob.glob(os.path.join(directory, pattern), recursive=True):
                content = extract_text_from_ppt(path)
                if not content or len(content.strip()) < 50:
                    continue
                title = os.path.splitext(os.path.basename(path))[0]
                item = KnowledgeBase(
                    title=title,
                    content=content,
                    category=category,
                    created_by=created_by,
                    is_active=True,
                )
                session.add(item)
                count += 1
        session.commit()
    finally:
        session.close()
    return count


def main() -> None:
    directory = os.environ.get("KB_PPT_DIR", r"D:\\Data_set")
    category = os.environ.get("KB_CATEGORY", "academic_notes")
    total = import_ppt_dir(directory, category)
    print(f"Imported {total} PPT/PPTX files from {directory}")


if __name__ == "__main__":
    main()


