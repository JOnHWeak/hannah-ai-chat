#!/usr/bin/env python3
"""
Script để import toàn bộ dữ liệu từ folder Data_set vào knowledge base
Hỗ trợ cả file PowerPoint (.ppt/.pptx) và có thể mở rộng cho các định dạng khác
"""

import glob
import os
import re
from typing import Optional, Tuple, List
from pathlib import Path

from database import SessionLocal
from models import KnowledgeBase


def extract_text_from_ppt(path: str) -> str:
    """Extract text content from PowerPoint files"""
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise RuntimeError(
            "Missing dependency python-pptx. Install with: pip install python-pptx"
        ) from exc

    try:
        presentation = Presentation(path)
        texts = []
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in paragraph.runs).strip()
                        if line:
                            slide_texts.append(line)
            
            if slide_texts:
                texts.append(f"--- Slide {slide_num} ---")
                texts.extend(slide_texts)
                texts.append("")  # Empty line between slides
        
        return "\n".join(texts)
    
    except Exception as e:
        # Fallback: tạo content từ filename nếu không đọc được file
        filename = os.path.basename(path)
        title = os.path.splitext(filename)[0]
        
        # Tạo content cơ bản từ tên file
        content = f"Document: {title}\n\n"
        content += f"This document contains information about {title}.\n"
        content += f"File location: {path}\n\n"
        content += "Note: This file could not be fully processed due to format compatibility issues.\n"
        content += "The content may need to be manually reviewed or converted to a newer format.\n"
        
        return content


def categorize_by_filename(filename: str) -> str:
    """Tự động phân loại dựa trên tên file"""
    filename_lower = filename.lower()
    
    # CSI106 courses
    if "csi106" in filename_lower:
        return "csi106_computer_science"
    
    # Database topics
    if any(keyword in filename_lower for keyword in ["database", "dbi", "sql"]):
        return "database_systems"
    
    # Data Structures and Algorithms
    if any(keyword in filename_lower for keyword in ["dsa", "data structure", "algorithm"]):
        return "data_structures_algorithms"
    
    # Specific DSA topics
    if any(keyword in filename_lower for keyword in ["sorting", "sort"]):
        return "sorting_algorithms"
    
    if any(keyword in filename_lower for keyword in ["tree", "trees"]):
        return "tree_data_structures"
    
    if any(keyword in filename_lower for keyword in ["graph", "graphs"]):
        return "graph_algorithms"
    
    if any(keyword in filename_lower for keyword in ["hash", "hashing"]):
        return "hashing_data_structures"
    
    if any(keyword in filename_lower for keyword in ["stack", "queue", "list"]):
        return "linear_data_structures"
    
    if any(keyword in filename_lower for keyword in ["recursion", "recursive"]):
        return "recursion_algorithms"
    
    if any(keyword in filename_lower for keyword in ["text", "string", "processing"]):
        return "text_processing"
    
    if any(keyword in filename_lower for keyword in ["complexity", "analysis"]):
        return "algorithm_analysis"
    
    # Chapter-based classification
    if "chapter" in filename_lower:
        return "course_chapters"
    
    # Default category
    return "academic_materials"


def clean_content(content: str) -> str:
    """Clean và format content để tối ưu cho AI learning"""
    if not content:
        return ""
    
    # Remove excessive whitespace
    content = re.sub(r'\n\s*\n\s*\n+', '\n\n', content)
    
    # Remove slide separators that are too repetitive
    content = re.sub(r'--- Slide \d+ ---\s*\n\s*--- Slide \d+ ---', '', content)
    
    # Clean up bullet points and formatting
    content = re.sub(r'^\s*[-•*]\s*', '• ', content, flags=re.MULTILINE)
    
    return content.strip()


def import_dataset_directory(
    directory: str,
    created_by: str = "dataset_import",
    min_content_length: int = 100
) -> Tuple[int, List[dict]]:
    """
    Import toàn bộ dataset từ directory vào knowledge base
    
    Returns:
        Tuple[int, List[dict]]: (total_imported, import_details)
    """
    session = SessionLocal()
    count = 0
    import_details = []
    
    try:
        # Tìm tất cả file PowerPoint
        patterns = ["**/*.ppt", "**/*.pptx"]
        
        for pattern in patterns:
            for path in glob.glob(os.path.join(directory, pattern), recursive=True):
                try:
                    # Extract content
                    content = extract_text_from_ppt(path)
                    content = clean_content(content)
                    
                    # Skip nếu content quá ngắn
                    if len(content.strip()) < min_content_length:
                        continue
                    
                    # Generate title và category
                    filename = os.path.basename(path)
                    title = os.path.splitext(filename)[0]
                    category = categorize_by_filename(filename)
                    
                    # Tạo knowledge base entry
                    item = KnowledgeBase(
                        title=title,
                        content=content,
                        category=category,
                        created_by=created_by,
                        is_active=True,
                    )
                    
                    session.add(item)
                    count += 1
                    
                    # Track import details
                    import_details.append({
                        "filename": filename,
                        "title": title,
                        "category": category,
                        "content_length": len(content),
                        "path": path
                    })
                    
                    print(f"✓ Imported: {filename} -> {category}")
                    
                except Exception as e:
                    print(f"✗ Failed to import {path}: {str(e)}")
                    continue
        
        session.commit()
        print(f"\n🎉 Successfully imported {count} files from {directory}")
        
    except Exception as e:
        session.rollback()
        print(f"❌ Import failed: {str(e)}")
        raise
    finally:
        session.close()
    
    return count, import_details


def print_import_summary(import_details: List[dict]):
    """In summary về việc import"""
    if not import_details:
        return
    
    print("\n" + "="*60)
    print("📊 IMPORT SUMMARY")
    print("="*60)
    
    # Group by category
    categories = {}
    for item in import_details:
        cat = item["category"]
        if cat not in categories:
            categories[cat] = []
        categories[cat].append(item)
    
    for category, items in categories.items():
        print(f"\n📁 {category.upper()}: {len(items)} files")
        for item in items:
            print(f"   • {item['filename']} ({item['content_length']} chars)")
    
    total_chars = sum(item["content_length"] for item in import_details)
    print(f"\n📈 Total: {len(import_details)} files, {total_chars:,} characters")


def main():
    """Main function để chạy import"""
    # Có thể override bằng environment variables
    dataset_dir = os.environ.get("DATASET_DIR", r"D:\Data_set")
    created_by = os.environ.get("IMPORT_CREATED_BY", "dataset_import")
    min_length = int(os.environ.get("MIN_CONTENT_LENGTH", "100"))
    
    print("🚀 Starting dataset import...")
    print(f"📂 Dataset directory: {dataset_dir}")
    print(f"👤 Created by: {created_by}")
    print(f"📏 Min content length: {min_length}")
    
    if not os.path.exists(dataset_dir):
        print(f"❌ Dataset directory not found: {dataset_dir}")
        return
    
    try:
        count, details = import_dataset_directory(
            directory=dataset_dir,
            created_by=created_by,
            min_content_length=min_length
        )
        
        print_import_summary(details)
        
        if count > 0:
            print(f"\n✅ Import completed successfully!")
            print(f"📚 {count} files added to knowledge base")
            print("\n💡 Next steps:")
            print("   1. Run: python scripts/ingest/index_kb_to_es.py")
            print("   2. Run: python scripts/sft/kb_to_sft.py")
            print("   3. Start training with: python scripts/train/train_lora_unsloth.py")
        else:
            print("⚠️  No files were imported. Check your dataset directory.")
            
    except Exception as e:
        print(f"❌ Import failed: {str(e)}")


if __name__ == "__main__":
    main()
