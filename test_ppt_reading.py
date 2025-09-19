#!/usr/bin/env python3
"""
Test script để kiểm tra việc đọc PowerPoint files từ dataset
"""

import os
import glob
from pptx import Presentation


def test_ppt_reading():
    """Test đọc một vài file PowerPoint từ dataset"""
    dataset_dir = r"D:\Data_set"
    
    if not os.path.exists(dataset_dir):
        print(f"❌ Dataset directory not found: {dataset_dir}")
        return
    
    print(f"🔍 Scanning dataset directory: {dataset_dir}")
    
    # Tìm các file PowerPoint
    patterns = ["**/*.ppt", "**/*.pptx"]
    ppt_files = []
    
    for pattern in patterns:
        ppt_files.extend(glob.glob(os.path.join(dataset_dir, pattern), recursive=True))
    
    print(f"📁 Found {len(ppt_files)} PowerPoint files")
    
    if not ppt_files:
        print("❌ No PowerPoint files found!")
        return
    
    # Test đọc một vài file đầu tiên
    test_files = ppt_files[:3]  # Test 3 files đầu tiên
    
    for i, file_path in enumerate(test_files, 1):
        print(f"\n📄 Test {i}: {os.path.basename(file_path)}")
        
        try:
            presentation = Presentation(file_path)
            print(f"   📊 Slides: {len(presentation.slides)}")
            
            # Đếm text content
            total_text_length = 0
            slides_with_text = 0
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            line = "".join(run.text for run in paragraph.runs).strip()
                            if line:
                                slide_text += line + " "
                
                if slide_text.strip():
                    slides_with_text += 1
                    total_text_length += len(slide_text.strip())
            
            print(f"   📝 Slides with text: {slides_with_text}/{len(presentation.slides)}")
            print(f"   📏 Total text length: {total_text_length} characters")
            
            if total_text_length > 0:
                print(f"   ✅ File readable and has content")
            else:
                print(f"   ⚠️  File readable but no text content")
                
        except Exception as e:
            print(f"   ❌ Error reading file: {str(e)}")
    
    print(f"\n🎉 PowerPoint reading test completed!")
    print(f"📊 Summary: {len(test_files)} files tested from {len(ppt_files)} total files")


if __name__ == "__main__":
    test_ppt_reading()
